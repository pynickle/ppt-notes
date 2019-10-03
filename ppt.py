import zipfile
import os
import shutil

from pptx import Presentation
from pptx.util import Pt
from bs4 import BeautifulSoup

from animation import animation, before, after

filename = "notes.txt"
rows_in_page = 10

with open(filename, "r", encoding="utf-8") as f:
    data = f.read()

data = data.split("\n\n")

ppt = Presentation()
width = ppt.slide_width
height = ppt.slide_height
pages_rows = []

for word in data:
    rows = word.split("\n")
    page = len(rows) // rows_in_page
    for i in range(page+1):
        slide = ppt.slides.add_slide(ppt.slide_layouts[6])

        textbox = slide.shapes.add_textbox(500000, 500000, width-1000000, height-1000000).text_frame
        for index in range(rows_in_page):
            if index == 0:
                textbox_para = textbox.paragraphs[0]
            else:
                textbox_para = textbox.add_paragraph()
            try:
                textbox_para.text = rows[rows_in_page*i + index]
            except IndexError:
                pages_rows.append(index)
                break
            else:
                textbox_para.font.name = "monospace"
                textbox_para.font.size = Pt(36)
        else:
            pages_rows.append(index + 1)

ppt.save("test.pptx")

f = zipfile.ZipFile("test.pptx",'r')
for file in f.namelist():
    f.extract(file,"./temp/")

index = 0
for i in os.listdir("./temp/ppt/slides"):
    if i == '_rels':
        break
    with open('./temp/ppt/slides/' + i, "r", encoding="utf-8") as f:
        soup = f.read()
    res = before
    res += "".join(animation[:pages_rows[index]])
    res += after
    soup = soup.replace("</p:clrMapOvr></p:sld>", "</p:clrMapOvr>" + res + "</p:sld>")
    with open("./temp/ppt/slides/" + i, "w", encoding="utf-8") as f:
        f.write(soup)
    index += 1

os.chdir("./temp")
with zipfile.ZipFile('../res.pptx', 'w') as f:
    for root, dirs, files in os.walk("."):
        for file in files:
            f.write(root + "/" + file) 

os.chdir("../")
shutil.rmtree("temp")
os.remove("test.pptx")

print("ppt已生成至res.pptx！")
